import os
import time
import hashlib
import pickle
from pathlib import Path
from typing import List
from langchain.schema import Document
from langchain_community.document_loaders import PyMuPDFLoader, UnstructuredFileLoader
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from pptx import Presentation  # type: ignore
from bs4 import BeautifulSoup
import requests
import argparse
import logging


# ========================
# üîß Logging setup
# ========================
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# ========================
# üîß Dynamic Paths (Repo-relative)
# ========================
# BASE_DIR ‚Üí Raggers/utils/
BASE_DIR = Path(__file__).resolve().parents[2]
# PROJECT_ROOT ‚Üí Chatbot/
PROJECT_ROOT = BASE_DIR.parent.parent

# Important folders (auto-adjust when repo is cloned anywhere)
HASH_STORE_PATH = BASE_DIR / "indexed_hashes.pkl"
INDEX_PATH = PROJECT_ROOT / "combined_faiss_index"
DEFAULT_DOC_FOLDER = PROJECT_ROOT / "Raggers" / "backend_rag_data"

# ========================
# üîß Constants
# ========================
SUPPORTED_EXTENSIONS = [".pdf", ".txt", ".md", ".csv", ".docx", ".ppt", ".pptx"]
MIN_TOKENS = 20
CHUNK_SIZE = 500
CHUNK_OVERLAP = 50


# ========================
# üîë Load or initialize hash store
# ========================
if HASH_STORE_PATH.exists():
    with open(HASH_STORE_PATH, "rb") as f:
        indexed_hashes = pickle.load(f)
else:
    indexed_hashes = set()


# ========================
# üîç Helper Functions
# ========================
def hash_content(text: str) -> str:
    return hashlib.md5(text.encode()).hexdigest()


def load_ppt_file(path: str) -> List[Document]:
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return [Document(page_content=text, metadata={"source": Path(path).name, "ingested_by": "backend"})]


def load_new_files(folder: Path, processed: set) -> List[Document]:
    docs = []
    for file in folder.glob("*"):
        ext = file.suffix.lower()
        if ext in SUPPORTED_EXTENSIONS and file.name not in processed:
            try:
                if ext == ".pdf":
                    loader = PyMuPDFLoader(str(file))
                    pages = loader.load()
                elif ext in [".ppt", ".pptx"]:
                    pages = load_ppt_file(str(file))
                else:
                    loader = UnstructuredFileLoader(str(file))
                    pages = loader.load()

                for i, doc in enumerate(pages):
                    doc.metadata["source"] = file.name
                    doc.metadata["page"] = i + 1
                    doc.metadata["ingested_by"] = "backend"

                docs.extend(pages)
                processed.add(file.name)
                logger.info(f"üìÑ Loaded {len(pages)} pages from {file.name}")
            except Exception as e:
                logger.error(f"‚ùå Failed to load {file.name}: {e}")
    return docs


def load_web(urls: List[str], url_cache: dict) -> List[Document]:
    docs = []
    for url in urls:
        try:
            response = requests.get(url, timeout=10)
            soup = BeautifulSoup(response.text, "html.parser")
            for tag in soup(["script", "style", "nav", "footer", "header", "noscript"]):
                tag.decompose()
            text = soup.get_text(separator="\n")
            cleaned = "\n".join([line.strip() for line in text.splitlines() if line.strip()])
            hash_val = hash_content(cleaned)
            if url_cache.get(url) == hash_val:
                logger.info(f"üîÑ No change in {url}, skipping...")
                continue
            url_cache[url] = hash_val
            doc = Document(page_content=cleaned, metadata={"source": url, "ingested_by": "backend"})
            docs.append(doc)
        except Exception as e:
            logger.error(f"‚ùå Failed to scrape {url}: {e}")
    return docs


def chunk_documents(docs: List[Document]) -> List[Document]:
    splitter = RecursiveCharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)
    chunks = splitter.split_documents(docs)
    filtered_chunks = []
    for i, chunk in enumerate(chunks):
        if len(chunk.page_content.strip().split()) >= MIN_TOKENS:
            chunk.metadata["chunk_index"] = i
            chunk.metadata["rag_snippet"] = chunk.page_content
            filtered_chunks.append(chunk)
    return filtered_chunks


def deduplicate_chunks(chunks: List[Document]) -> List[Document]:
    logger.warning("‚ö†Ô∏è Deduplication temporarily disabled ‚Äî all chunks will be indexed.")
    return chunks


def update_index(chunks: List[Document], index_path=INDEX_PATH):
    logger.info(f"üóÇÔ∏è Updating FAISS index at: {index_path}")
    embedder = HuggingFaceEmbeddings(model_name="BAAI/bge-small-en")

    if Path(index_path).exists():
        index = FAISS.load_local(index_path, embedder, allow_dangerous_deserialization=True)
        index.add_documents(chunks)
    else:
        os.makedirs(index_path, exist_ok=True)
        index = FAISS.from_documents(chunks, embedder)

    index.save_local(index_path)
    logger.info(f"‚úÖ Index updated and saved to '{index_path}'")
    logger.info(f"üìä FAISS now contains {len(index.docstore._dict)} documents")

    with open(HASH_STORE_PATH, "wb") as f:
        pickle.dump(indexed_hashes, f)


# ========================
# üöÄ Main Ingestion Function
# ========================
def run_background_ingestion(pdf_dir: Path = DEFAULT_DOC_FOLDER, urls: List[str] = None,
                             index_path=INDEX_PATH, benchmark=False):
    if urls is None:
        urls = []
    start = time.time()

    pdf_dir = Path(pdf_dir)
    logger.info(f"üöÄ Ingesting from folder: {pdf_dir}")

    if not pdf_dir.exists():
        logger.error(f"‚ùå Folder does not exist: {pdf_dir}")
        return

    processed_files = set()
    url_cache = {}

    new_file_docs = load_new_files(pdf_dir, processed_files)
    new_web_docs = load_web(urls, url_cache)

    all_docs = new_file_docs + new_web_docs
    if not all_docs:
        logger.warning(f"‚ö†Ô∏è No new documents found in {pdf_dir}")
        return

    chunks = chunk_documents(all_docs)
    chunks = deduplicate_chunks(chunks)

    if chunks:
        logger.info(f"‚úÖ {len(chunks)} chunks to index.")
        update_index(chunks, index_path)
    else:
        logger.warning("‚ùå No valid chunks to index.")

    if benchmark:
        logger.info(f"‚è±Ô∏è Ingestion completed in {round(time.time() - start, 2)}s")


# ========================
# üß© CLI Entry Point
# ========================
if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Backend ingestion script for RAG")
    parser.add_argument("--folder", type=str, default=str(DEFAULT_DOC_FOLDER),
                        help="Folder containing docs to ingest")
    parser.add_argument("--update", action="store_true", help="Update the existing FAISS index")
    parser.add_argument("--benchmark", action="store_true", help="Measure ingestion time")
    parser.add_argument("--index", type=str, default=str(INDEX_PATH), help="Path to FAISS index directory")

    args = parser.parse_args()

    run_background_ingestion(
        pdf_dir=args.folder,
        urls=[],
        index_path=args.index,
        benchmark=args.benchmark
    )
