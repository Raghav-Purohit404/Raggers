import os
import sys
from pathlib import Path
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.document_loaders import PyMuPDFLoader, UnstructuredFileLoader
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter

# ==============================
# ⚙️ Path Configuration
# ==============================
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(os.path.dirname(BASE_DIR))

DATA_FOLDER = os.path.join(PROJECT_ROOT, "Raggers", "backend_rag_data")
INDEX_PATH = os.path.join(PROJECT_ROOT, "Raggers", "combined_faiss_index")

SUPPORTED_EXTENSIONS = [".pdf", ".txt", ".md", ".csv", ".docx", ".ppt", ".pptx"]

# ==============================
# 📄 File Loading Utilities
# ==============================
def load_ppt_file(path: str):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def load_documents(folder: str):
    docs = []
    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    for file in Path(folder).glob("*"):
        ext = file.suffix.lower()
        if ext in SUPPORTED_EXTENSIONS:
            try:
                if ext == ".pdf":
                    loader = PyMuPDFLoader(str(file))
                    pages = loader.load()
                elif ext in [".ppt", ".pptx"]:
                    text = load_ppt_file(str(file))
                    pages = [{"page_content": text}]
                else:
                    loader = UnstructuredFileLoader(str(file))
                    pages = loader.load()
                for doc in pages:
                    chunks = splitter.split_text(doc["page_content"])
                    docs.extend(chunks)
            except Exception as e:
                print(f"❌ Failed to load {file.name}: {e}")
    return docs

# ==============================
# 🧠 Rebuild FAISS
# ==============================
def rebuild_faiss():
    if not os.path.exists(DATA_FOLDER):
        print(f"❌ Folder does not exist: {DATA_FOLDER}")
        return

    docs = load_documents(DATA_FOLDER)
    if not docs:
        print("⚠️ No documents found to index.")
        return

    embedder = HuggingFaceEmbeddings(model_name="BAAI/bge-small-en")

    if os.path.exists(INDEX_PATH):
        import shutil
        shutil.rmtree(INDEX_PATH)

    os.makedirs(INDEX_PATH, exist_ok=True)
    index = FAISS.from_texts(docs, embedder)
    index.save_local(INDEX_PATH)
    print(f"✅ FAISS index rebuilt successfully with {len(docs)} chunks.")
    print(f"📁 Index saved to: {INDEX_PATH}")

# ==============================
# 🚀 Run
# ==============================
if __name__ == "__main__":
    rebuild_faiss()
